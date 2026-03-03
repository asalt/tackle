from __future__ import annotations

import os
import re
import textwrap
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence, Tuple

import pandas as pd

from .exporter import _collect_tsvs, _read_tsv
from .utils import _get_logger

logger = _get_logger(__name__)


@dataclass(frozen=True)
class TableAsset:
    title: str
    source_relpath: str
    source_path: str
    table_rows: int
    table_cols: int
    png_path: Optional[str]
    pdf_path: Optional[str]


@dataclass(frozen=True)
class DeckImageAsset:
    title: str
    source_relpath: str
    source_path: str
    png_path: str
    category: str


def _sanitize_filename(name: str) -> str:
    safe = re.sub(r"[^\w.\-]+", "_", str(name)).strip("._-")
    return safe or "table"


def _readable_title_from_rel(relpath: str) -> str:
    base = os.path.basename(relpath)
    base = re.sub(r"\.tsv$", "", base, flags=re.IGNORECASE)
    # Prefer extracting everything after 'group' if present (matches exporter/make-xls behavior)
    m = re.search(r"(?<=group)[_]?(.*)$", base)
    label = m.group(1) if m else base
    label = label.replace("__", " ").replace("_", " ").strip()
    label = re.sub(r"\s+", " ", label)
    return label or base


def _readable_title_from_png_rel(relpath: str, *, category: str) -> str:
    p = Path(relpath)
    stem = re.sub(r"[_\-]+", " ", p.stem).strip()
    parent = re.sub(r"[_\-]+", " ", p.parent.name).strip() if p.parent.name else ""
    category_label = {
        "volcano": "Volcano Plot",
        "pca": "PCA Plot",
        "metrics": "Metrics Plot",
        "cluster": "Cluster Plot",
        "topdiff-cluster": "Topdiff Cluster Plot",
    }.get(category, category.title())

    if parent and parent.lower() not in {"volcano", "pca", "metrics", "cluster", "clustermap", "topdiff"}:
        return f"{category_label}: {parent} / {stem or p.name}"
    return f"{category_label}: {stem or p.name}"


def _classify_plot_png(relpath: str) -> Optional[str]:
    rel_lower = relpath.lower().replace("\\", "/")
    # topdiff-specific cluster images should be categorized first.
    if "topdiff" in rel_lower and ("cluster" in rel_lower or "clustermap" in rel_lower):
        return "topdiff-cluster"
    if "volcano" in rel_lower:
        return "volcano"
    if "pca" in rel_lower or "pcaplot" in rel_lower:
        return "pca"
    if "metrics" in rel_lower:
        return "metrics"
    if "cluster" in rel_lower or "clustermap" in rel_lower:
        return "cluster"
    return None


def _is_under(path: Path, parent: Path) -> bool:
    try:
        path.resolve().relative_to(parent.resolve())
        return True
    except Exception:
        return False


def collect_plot_image_assets(
    *,
    base_dir: str,
    filter_contains: Sequence[str] = (),
    include_kinds: Sequence[str] = ("volcano", "pca", "metrics", "cluster", "topdiff-cluster"),
    exclude_dirs: Sequence[str] = (),
) -> List[DeckImageAsset]:
    """
    Discover existing PNG plot assets (volcano/pca/metrics/cluster/topdiff-cluster)
    to include in slide decks.
    """
    root = Path(base_dir).resolve()
    include = {str(x).strip().lower() for x in include_kinds if str(x).strip()}
    needles = [s for s in filter_contains if s]
    excluded = [Path(p).resolve() for p in exclude_dirs if str(p).strip()]
    # Always exclude report/deck outputs to avoid re-ingesting generated assets.
    excluded.append((root / "report" / "deck").resolve())

    assets: List[DeckImageAsset] = []
    for png in root.rglob("*.png"):
        png_resolved = png.resolve()
        if any(_is_under(png_resolved, ex) for ex in excluded):
            continue

        rel = png_resolved.relative_to(root)
        rel_str = str(rel)
        if needles and not any(n in rel_str for n in needles):
            continue

        category = _classify_plot_png(rel_str)
        if not category or category not in include:
            continue

        assets.append(
            DeckImageAsset(
                title=_readable_title_from_png_rel(rel_str, category=category),
                source_relpath=rel_str,
                source_path=str(png_resolved),
                png_path=str(png_resolved),
                category=category,
            )
        )

    assets.sort(key=lambda a: a.source_relpath)
    return assets


def _pick_p_column(columns: Sequence[str]) -> Optional[str]:
    cols = set(str(c) for c in columns)
    if "pAdj" in cols:
        return "pAdj"
    if "pValue" in cols:
        return "pValue"
    return None


def _summarize_volcano_top_hits(
    df: pd.DataFrame,
    *,
    topn: int = 15,
    direction: str = "both",
    p_cutoff: Optional[float] = 0.05,
    fc_cutoff: Optional[float] = None,
) -> Tuple[pd.DataFrame, Dict[str, Any]]:
    if direction not in ("both", "up", "down"):
        raise ValueError(f"direction must be one of both/up/down, got {direction}")
    if topn < 1:
        raise ValueError("topn must be >= 1")

    df = df.copy()
    df.columns = [str(c) for c in df.columns]

    if "GeneID" not in df.columns:
        # If GeneID is an index (older exports), promote it.
        if df.index.name and str(df.index.name).lower() == "geneid":
            df = df.reset_index()
        else:
            raise ValueError("Volcano table is missing required GeneID column")

    p_col = _pick_p_column(df.columns)
    if p_col is None:
        raise ValueError("Volcano table must contain pAdj or pValue")
    if "log2_FC" not in df.columns:
        raise ValueError("Volcano table must contain log2_FC")

    # Basic cleanup / typing
    df["GeneID"] = df["GeneID"].astype(str)
    df["log2_FC"] = pd.to_numeric(df["log2_FC"], errors="coerce")
    df[p_col] = pd.to_numeric(df[p_col], errors="coerce")

    # Optional filtering
    if p_cutoff is not None and float(p_cutoff) > 0:
        df = df.loc[df[p_col].notna() & (df[p_col] <= float(p_cutoff))]
    if fc_cutoff is not None:
        df = df.loc[df["log2_FC"].notna() & (df["log2_FC"].abs() >= float(fc_cutoff))]

    def _pick_dir(_df: pd.DataFrame, *, label: str, mask) -> pd.DataFrame:
        sub = _df.loc[mask].copy()
        if sub.empty:
            return sub
        # Prefer p-value sorting; tie-break by abs(fc)
        sub["_abs_fc"] = sub["log2_FC"].abs()
        sub = sub.sort_values(by=[p_col, "_abs_fc"], ascending=[True, False], kind="stable").head(int(topn))
        sub = sub.drop(columns=["_abs_fc"], errors="ignore")
        sub.insert(0, "Direction", label)
        return sub

    if direction == "up":
        picked = _pick_dir(df, label="Up", mask=df["log2_FC"] > 0)
    elif direction == "down":
        picked = _pick_dir(df, label="Down", mask=df["log2_FC"] < 0)
    else:
        up = _pick_dir(df, label="Up", mask=df["log2_FC"] > 0)
        down = _pick_dir(df, label="Down", mask=df["log2_FC"] < 0)
        picked = pd.concat([up, down], axis=0, ignore_index=True)

    # Columns to show (keep it slide-friendly)
    show_cols: List[str] = []
    if "Direction" in picked.columns:
        show_cols.append("Direction")
    show_cols.append("GeneSymbol" if "GeneSymbol" in picked.columns else "GeneID")
    if "GeneSymbol" in picked.columns and "GeneID" in picked.columns:
        show_cols.append("GeneID")
    show_cols.extend(["log2_FC", p_col])
    if "GeneDescription" in picked.columns:
        show_cols.append("GeneDescription")
    elif "Description" in picked.columns:
        show_cols.append("Description")
    if "FunCats" in picked.columns:
        show_cols.append("FunCats")

    picked = picked[[c for c in show_cols if c in picked.columns]].copy()

    meta = {
        "p_col": p_col,
        "p_cutoff": p_cutoff,
        "fc_cutoff": fc_cutoff,
        "topn": topn,
        "direction": direction,
        "total_rows": int(df.shape[0]),
        "picked_rows": int(picked.shape[0]),
    }
    return picked, meta


def _format_table_cell(value: Any, *, col: str) -> str:
    if value is None:
        return ""
    try:
        if pd.isna(value):
            return ""
    except Exception:
        pass

    if col in ("pAdj", "pValue"):
        try:
            return f"{float(value):.2E}"
        except Exception:
            return str(value)
    if col in ("log2_FC", "signedlogP", "t"):
        try:
            return f"{float(value):+.2f}"
        except Exception:
            return str(value)

    s = str(value)
    s = s.replace("\r\n", "\n").replace("\r", "\n").replace("\t", " ")
    s = re.sub(r"\s+", " ", s).strip()
    if col in ("GeneDescription", "Description"):
        if len(s) > 100:
            s = s[:97].rstrip() + "..."
    if col == "FunCats" and len(s) > 80:
        s = s[:77].rstrip() + "..."
    return s


def _render_table_figure(
    df: pd.DataFrame,
    *,
    title: str,
    subtitle: Optional[str] = None,
) -> Any:
    # Avoid GUI backends in headless environments.
    import matplotlib

    try:
        matplotlib.use("Agg", force=True)
    except Exception:
        pass
    import matplotlib.pyplot as plt

    nrows, ncols = df.shape
    # Widescreen-ish defaults; table height grows with row count.
    width_in = min(13.33, max(9.0, 1.25 * ncols + 4.0))
    height_in = min(7.5, max(4.0, 0.28 * max(1, nrows) + 2.2))

    fig = plt.figure(figsize=(width_in, height_in))

    # Title / subtitle in figure coordinates
    fig.text(
        0.5,
        0.97,
        title,
        ha="center",
        va="top",
        fontsize=16,
        fontweight="bold",
    )
    if subtitle:
        fig.text(0.5, 0.93, subtitle, ha="center", va="top", fontsize=10, color="#444444")

    ax = fig.add_axes([0.03, 0.05, 0.94, 0.85])
    ax.axis("off")

    headers = [str(c) for c in df.columns]
    cell_text = [
        [_format_table_cell(v, col=headers[j]) for j, v in enumerate(row)]
        for row in df.itertuples(index=False, name=None)
    ]

    table = ax.table(
        cellText=cell_text,
        colLabels=headers,
        cellLoc="left",
        colLoc="left",
        loc="upper center",
    )
    table.auto_set_font_size(False)
    # Shrink fonts as the table grows.
    base_size = 11
    font_size = max(6, int(base_size - (nrows / 18) - (ncols / 14)))
    table.set_fontsize(font_size)

    header_bg = "#F2F2F2"
    row_bg_a = "#FFFFFF"
    row_bg_b = "#FAFAFA"
    edge = "#DDDDDD"

    numeric_cols = {c for c in headers if c in ("pAdj", "pValue", "log2_FC", "signedlogP", "t")}

    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor(edge)
        cell.set_linewidth(0.6)
        if r == 0:
            cell.set_facecolor(header_bg)
            cell.get_text().set_weight("bold")
            cell.get_text().set_color("#222222")
            cell.get_text().set_ha("left")
        else:
            cell.set_facecolor(row_bg_b if (r % 2 == 0) else row_bg_a)
            if headers[c] in numeric_cols:
                cell.get_text().set_ha("right")
            else:
                cell.get_text().set_ha("left")

    # Column widths based on content length (approx), with emphasis on description.
    col_lens: List[int] = []
    for col in headers:
        try:
            max_len = max([len(col), *df[col].astype(str).map(len).tolist()])
        except Exception:
            max_len = len(col)
        col_lens.append(max_len)

    weights = []
    for col, ln in zip(headers, col_lens):
        w = max(8, min(40, int(ln)))
        if col in ("GeneDescription", "Description"):
            w = int(w * 1.7)
        elif col in ("FunCats",):
            w = int(w * 1.3)
        weights.append(w)

    tot = float(sum(weights)) if weights else 1.0
    col_widths = [w / tot for w in weights]

    for col_idx, w in enumerate(col_widths):
        for row_idx in range(0, nrows + 1):
            table[(row_idx, col_idx)].set_width(w)

    # Slight vertical scaling for readability.
    table.scale(1.0, 1.18)
    return fig


def build_volcano_table_assets(
    *,
    base_dir: str,
    out_dir: str,
    filter_contains: Sequence[str] = (),
    topn: int = 15,
    direction: str = "both",
    p_cutoff: Optional[float] = 0.05,
    fc_cutoff: Optional[float] = None,
    formats: Sequence[str] = (".png", ".pdf"),
    dpi: int = 300,
    tsv_engine: str = "auto",
    pandas_low_memory: bool = False,
    force: bool = False,
) -> List[TableAsset]:
    """
    Generate slide-ready table images summarizing volcano TSVs under an analysis directory.

    Returns a list of TableAsset entries including written paths.
    """
    root = Path(base_dir)
    out_root = Path(out_dir)
    out_root.mkdir(parents=True, exist_ok=True)

    patterns = [("volcano", "*.tsv")]
    targets = _collect_tsvs(root, patterns)

    if filter_contains:
        needles = [s for s in filter_contains if s]
        if needles:
            targets = [(rel, p) for (rel, p) in targets if any(n in str(rel) for n in needles)]

    if not targets:
        raise FileNotFoundError(f"No volcano TSVs found under {base_dir}")

    formats = tuple(f if f.startswith(".") else f".{f}" for f in formats)
    supported = {".png", ".pdf"}
    if any(f not in supported for f in formats):
        bad = ", ".join(sorted(set(f for f in formats if f not in supported)))
        raise ValueError(f"Unsupported formats: {bad} (supported: png,pdf)")

    assets: List[TableAsset] = []
    for rel, path in targets:
        rel_str = str(rel)
        title = _readable_title_from_rel(rel_str)

        try:
            df = _read_tsv(path, engine=tsv_engine, low_memory=pandas_low_memory)
            summary, meta = _summarize_volcano_top_hits(
                df,
                topn=topn,
                direction=direction,
                p_cutoff=p_cutoff,
                fc_cutoff=fc_cutoff,
            )
        except Exception as e:
            logger.warning("Deckgen: skipping %s (%s)", rel_str, e)
            continue

        if summary.empty:
            # Preserve a useful placeholder with a single row.
            summary = pd.DataFrame({"Note": ["No rows matched filters (or no hits in this direction)."]})

        subtitle_parts = []
        if meta.get("p_col"):
            subtitle_parts.append(f"{meta['p_col']}≤{meta.get('p_cutoff')}" if (meta.get("p_cutoff") and meta.get("p_cutoff") > 0) else str(meta["p_col"]))
        if meta.get("fc_cutoff") is not None:
            subtitle_parts.append(f"|log2_FC|≥{meta['fc_cutoff']}")
        subtitle_parts.append(f"top {meta.get('topn')} ({meta.get('direction')})")
        subtitle = " · ".join(str(x) for x in subtitle_parts if x)

        fig = _render_table_figure(summary, title=title, subtitle=subtitle)

        stem = _sanitize_filename(rel_str.replace(os.sep, "__").replace("/", "__"))
        stem = re.sub(r"\.tsv$", "", stem, flags=re.IGNORECASE)
        out_base = out_root / stem

        written_png = None
        written_pdf = None
        for fmt in formats:
            out_path = Path(str(out_base) + fmt)
            if out_path.exists() and not force:
                logger.info("Deckgen: skipping existing %s (use --force to overwrite)", out_path)
                if fmt == ".png":
                    written_png = str(out_path)
                elif fmt == ".pdf":
                    written_pdf = str(out_path)
                continue
            fig.savefig(out_path, dpi=dpi, bbox_inches="tight", pad_inches=0.15)
            if fmt == ".png":
                written_png = str(out_path)
            elif fmt == ".pdf":
                written_pdf = str(out_path)

        try:
            import matplotlib.pyplot as plt

            plt.close(fig)
        except Exception:
            pass

        assets.append(
            TableAsset(
                title=title,
                source_relpath=rel_str,
                source_path=str(path),
                table_rows=int(summary.shape[0]),
                table_cols=int(summary.shape[1]),
                png_path=written_png,
                pdf_path=written_pdf,
            )
        )

    assets.sort(key=lambda a: a.source_relpath)
    return assets


def build_pptx_deck(
    *,
    out_path: str,
    title: str,
    assets: Sequence[TableAsset | DeckImageAsset],
    subtitle: Optional[str] = None,
    widescreen: bool = True,
) -> str:
    """
    Create a PowerPoint (.pptx) with one slide per PNG asset.
    Requires `python-pptx`.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as e:
        raise RuntimeError(
            "python-pptx is required to generate slide decks. Install via: pip install 'tackle[slides]'"
        ) from e

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    def _add_title(slide, text: str):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), slide_w - Inches(1.0), Inches(0.7))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(28)
        run.font.bold = True

    def _add_subtitle(slide, text: str):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), slide_w - Inches(1.0), Inches(0.4))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(14)

    # Title slide
    s0 = prs.slides.add_slide(blank)
    _add_title(s0, title)
    if subtitle:
        _add_subtitle(s0, subtitle)
    else:
        _add_subtitle(s0, datetime.now().strftime("%Y-%m-%d %H:%M"))

    # Content slides
    top = Inches(1.2)
    left = Inches(0.5)
    max_w = slide_w - Inches(1.0)
    max_h = slide_h - Inches(1.6)

    for asset in assets:
        png_path = getattr(asset, "png_path", None)
        slide_title = getattr(asset, "title", None)
        if not png_path:
            continue
        if not slide_title:
            slide_title = Path(str(png_path)).stem
        slide = prs.slides.add_slide(blank)
        _add_title(slide, str(slide_title))

        pic = slide.shapes.add_picture(str(png_path), left, top)
        # Scale to fit box, preserving aspect.
        scale = min(max_w / pic.width, max_h / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(left + (max_w - pic.width) / 2)
        pic.top = int(top + (max_h - pic.height) / 2)

    prs.save(out)
    return str(out)
